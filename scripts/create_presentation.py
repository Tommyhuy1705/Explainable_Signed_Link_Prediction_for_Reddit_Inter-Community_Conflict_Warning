"""Create the final 12-slide PowerPoint presentation from report artifacts."""

from __future__ import annotations

from pathlib import Path
import sys

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]

INK = RGBColor(34, 34, 34)
MUTED = RGBColor(96, 104, 116)
RED = RGBColor(143, 29, 44)
BLUE = RGBColor(47, 93, 140)
GREEN = RGBColor(61, 139, 91)
PAPER = RGBColor(248, 247, 244)
LIGHT = RGBColor(232, 235, 239)


def _add_text(slide, text, left, top, width, height, *, size=20, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    return box


def _add_kicker(slide, text, number):
    _add_text(slide, f"{number:02d}", 0.55, 0.35, 0.35, 0.22, size=8, color=RED, bold=True)
    _add_text(slide, text.upper(), 0.95, 0.33, 3.2, 0.25, size=8, color=MUTED, bold=True)


def _add_title(slide, kicker, number, title, subtitle=None):
    _add_kicker(slide, kicker, number)
    _add_text(slide, title, 0.55, 0.7, 7.5, 0.72, size=25, color=INK, bold=True)
    if subtitle:
        _add_text(slide, subtitle, 0.58, 1.38, 7.2, 0.45, size=10.5, color=MUTED)


def _add_footer(slide, page):
    _add_text(slide, "UEH Social Media Data Analysis | Reddit signed-network final project", 0.55, 7.05, 6.2, 0.2, size=7.5, color=MUTED)
    _add_text(slide, str(page), 12.4, 7.05, 0.35, 0.2, size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_picture(slide, path, left, top, width=None, height=None):
    picture_path = ROOT / path
    if not picture_path.exists():
        raise FileNotFoundError(picture_path)
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    return slide.shapes.add_picture(str(picture_path), Inches(left), Inches(top), **kwargs)


def _add_bullet_list(slide, items, left, top, width, height, *, size=15):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = INK
        paragraph.font.name = "Aptos"
    return box


def _add_metric(slide, value, label, left, top, width=2.4, color=BLUE):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = LIGHT
    _add_text(slide, value, left + 0.16, top + 0.12, width - 0.3, 0.28, size=18, color=color, bold=True)
    _add_text(slide, label, left + 0.16, top + 0.48, width - 0.3, 0.24, size=8.5, color=MUTED)


def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    return slide


def create_deck(output_path: Path) -> Path:
    metrics = pd.read_csv(ROOT / "data/processed/phase3/phase3_model_metrics.csv")
    best = metrics.sort_values(["test_pr_auc", "test_f1"], ascending=False).iloc[0]
    robustness_path = ROOT / "data/processed/phase3/robustness_metrics.csv"
    robustness_note = "Robustness artifacts generated from k-core sensitivity probe."
    if robustness_path.exists():
        robustness = pd.read_csv(robustness_path)
        robustness_note = f"{len(robustness)} k-core sensitivity rows compare global and history-safe filtering."

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = _blank(prs)
    _add_text(slide, "Predicting Negative Cross-Community Hyperlinks on Reddit", 0.65, 0.75, 8.2, 1.2, size=31, color=INK, bold=True)
    _add_text(slide, "Temporal signed-network features for Reddit inter-community conflict warning", 0.72, 2.0, 6.9, 0.35, size=13.5, color=MUTED)
    _add_metric(slide, "858,488", "raw Kaggle/SNAP hyperlink records", 0.72, 3.0, color=BLUE)
    _add_metric(slide, "13+", "report-ready figures", 3.35, 3.0, color=GREEN)
    _add_metric(slide, "0.184", "best test PR-AUC", 5.98, 3.0, color=RED)
    _add_picture(slide, "reports/figures/community_network_sample.png", 8.0, 0.6, width=4.8)
    _add_footer(slide, 1)

    slide = _blank(prs)
    _add_title(slide, "Motivation", 2, "Negative community links are rare, structured, and worth predicting.", "The project turns Reddit hyperlinks into a signed directed network for early warning analysis.")
    _add_bullet_list(slide, [
        "Subreddits can amplify conflict through cross-community references.",
        "A negative hyperlink is a measurable proxy, not a claim of real-world harassment.",
        "The practical goal is to prioritize relationships that deserve monitoring.",
    ], 0.85, 2.1, 4.4, 2.4)
    _add_picture(slide, "reports/figures/monthly_negative_ratio.png", 5.45, 1.75, width=7.2)
    _add_footer(slide, 2)

    slide = _blank(prs)
    _add_title(slide, "Dataset", 3, "A Kaggle mirror keeps the data accessible; SNAP anchors provenance.", "Files are the Reddit Hyperlink title/body signed graph tables with 86 text-property features.")
    _add_metric(slide, "Kaggle", "Signed Graphs mirror", 0.75, 1.9, color=RED)
    _add_metric(slide, "2013-2017", "timestamp range", 3.35, 1.9, color=BLUE)
    _add_metric(slide, "9.58%", "raw negative-link share", 5.95, 1.9, color=GREEN)
    _add_picture(slide, "reports/figures/label_distribution.png", 0.75, 3.1, width=5.5)
    _add_bullet_list(slide, [
        "Raw schema: source, target, post id, timestamp, link sentiment, properties.",
        "Audit validates labels, timestamps, required columns, and exactly 86 properties.",
        "Original academic source: Kumar, Hamilton, Leskovec, Jurafsky, WWW 2018.",
    ], 6.65, 3.1, 5.7, 2.2, size=12)
    _add_footer(slide, 3)

    slide = _blank(prs)
    _add_title(slide, "Graph Formulation", 4, "The social-media task is a temporal signed link prediction problem.", "Nodes are subreddits; directed signed edges are cross-community hyperlinks.")
    _add_bullet_list(slide, [
        "History window: construct signed graph and pair-level features.",
        "Future window: label a pair negative when negative links outnumber positive/neutral links.",
        "This is relationship-level prediction, stricter than detecting one isolated negative edge.",
    ], 0.85, 1.9, 4.6, 2.4)
    _add_picture(slide, "reports/figures/community_pair_negative_heatmap.png", 5.6, 1.1, width=6.65)
    _add_footer(slide, 4)

    slide = _blank(prs)
    _add_title(slide, "Temporal Design", 5, "Strict history-label separation reduces leakage risk.", "Feature computation stops before the label window begins.")
    _add_bullet_list(slide, [
        "Train: history <= 2015-12-31, labels in 2016-H1.",
        "Validation: history <= 2016-06-30, labels in 2016-H2.",
        "Test: history <= 2016-12-31, labels through 2017-04-30.",
        "Decision thresholds are tuned on validation only.",
    ], 0.85, 1.8, 5.0, 2.7)
    _add_picture(slide, "reports/figures/monthly_negative_ratio.png", 5.75, 1.7, width=6.75)
    _add_footer(slide, 5)

    slide = _blank(prs)
    _add_title(slide, "Feature System", 6, "Hybrid features combine pair history, network position, community context, and text signals.", None)
    _add_bullet_list(slide, [
        "Pair history: positive/negative counts, negative ratio, reciprocity.",
        "Network: in/out degree, signed degree, PageRank, betweenness.",
        "Community: clustering, community size, negativity ratios, same-community flag.",
        "Balance: signed local-neighborhood triad counts.",
        "Text: 86 SNAP text-property features aggregated by source-target pair.",
    ], 0.9, 1.55, 5.6, 3.2, size=12.5)
    _add_picture(slide, "reports/figures/feature_importance_top20.png", 6.55, 1.15, width=5.95)
    _add_footer(slide, 6)

    slide = _blank(prs)
    _add_title(slide, "Network Evidence", 7, "Negative-link concentration is visible in community structure before modeling.", None)
    _add_picture(slide, "reports/figures/community_network_sample.png", 0.65, 1.15, width=7.1)
    _add_picture(slide, "reports/figures/community_negative_ratio.png", 8.05, 1.35, width=4.4)
    _add_footer(slide, 7)

    slide = _blank(prs)
    _add_title(slide, "Models", 8, "Baselines make the lift credible; ablations identify which signal family matters.", None)
    _add_bullet_list(slide, [
        "Baselines: dummy frequent, dummy prior, historical negative-ratio heuristic.",
        "ML models: Logistic Regression, Random Forest, XGBoost, LightGBM.",
        "Feature sets: history-only, text-only, graph-only, hybrid, and no-balance ablations.",
    ], 0.85, 1.75, 4.6, 2.1, size=12.5)
    _add_picture(slide, "reports/figures/model_comparison_pr_auc.png", 5.45, 0.95, width=7.25)
    _add_footer(slide, 8)

    slide = _blank(prs)
    _add_title(slide, "Main Result", 9, "Hybrid Logistic Regression is best by PR-AUC under strict temporal evaluation.", None)
    _add_metric(slide, f"{best['test_pr_auc']:.3f}", "test PR-AUC", 0.85, 1.75, color=RED)
    _add_metric(slide, f"{best['test_roc_auc']:.3f}", "test ROC-AUC", 3.45, 1.75, color=BLUE)
    _add_metric(slide, f"{best['test_f1']:.3f}", "negative-class F1", 6.05, 1.75, color=GREEN)
    _add_picture(slide, "reports/figures/best_confusion_matrix.png", 0.85, 3.05, width=4.55)
    _add_picture(slide, "reports/figures/precision_recall_curve.png", 5.95, 2.85, width=6.2)
    _add_footer(slide, 9)

    slide = _blank(prs)
    _add_title(slide, "Robustness", 10, "Sensitivity checks separate the main result from the graph-density choice.", robustness_note)
    _add_picture(slide, "reports/figures/robustness_kcore_pr_auc.png", 0.85, 1.55, width=5.5)
    _add_picture(slide, "reports/figures/threshold_tradeoff.png", 6.65, 1.55, width=5.6)
    _add_footer(slide, 10)

    slide = _blank(prs)
    _add_title(slide, "Limitations", 11, "The model predicts a proxy signal, not verified real-world conflict.", None)
    _add_bullet_list(slide, [
        "LINK_SENTIMENT is a derived label and can contain noise.",
        "K-core filtering improves tractability but narrows the evaluated subgraph.",
        "The current task predicts historically observed source-target pairs.",
        "The data is historical; platform behavior may differ today.",
    ], 0.95, 1.65, 5.4, 3.0, size=14)
    _add_picture(slide, "reports/figures/top_negative_sources.png", 6.6, 1.35, width=5.55)
    _add_footer(slide, 11)

    slide = _blank(prs)
    _add_title(slide, "Conclusion", 12, "Temporal signed-network features provide practical early-warning signals.", None)
    _add_bullet_list(slide, [
        "Graph/history features are the most stable signal family.",
        "Text properties help most when combined with network features.",
        "Structural-balance features add interpretability even when their lift is small.",
        "Future work: signed embeddings, temporal GNNs, SHAP, newer social-platform datasets.",
    ], 0.95, 1.75, 5.8, 3.1, size=14)
    _add_picture(slide, "reports/figures/model_comparison_pr_auc.png", 6.85, 1.35, width=5.35)
    _add_footer(slide, 12)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> int:
    output = ROOT / "docs" / "final_presentation.pptx"
    path = create_deck(output)
    print(path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
